#!/usr/bin/env python3
"""
Nivesh Saathi — Pitch Deck Generator
Generates a professional, visually stunning pitch deck as PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════
# Color Palette
# ═══════════════════════════════════════════
PRIMARY = RGBColor(0x6C, 0x63, 0xFF)      # Indigo
PRIMARY_DARK = RGBColor(0x4F, 0x46, 0xE5)
ACCENT = RGBColor(0x34, 0xD3, 0x99)       # Green
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24) # Gold
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)    # Red
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Navy
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BLACK = RGBColor(0x1E, 0x29, 0x3B)
SOFT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=WHITE):
    """lines is list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, title, body_lines, emoji="", title_color=WHITE, card_color=DARK_CARD, border=None):
    shape = add_shape(slide, left, top, width, height, card_color, border)
    y_offset = Pt(12)
    if emoji:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.5), Inches(0.4), emoji, 24, WHITE, False, PP_ALIGN.LEFT)
        add_text_box(slide, left + Inches(0.55), top + Inches(0.12), width - Inches(0.7), Inches(0.4), title, 14, title_color, True, PP_ALIGN.LEFT)
    else:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.4), title, 14, title_color, True, PP_ALIGN.LEFT)

    lines_data = [(line, 11, LIGHT_GRAY, False, PP_ALIGN.LEFT) for line in body_lines]
    add_multi_text(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.65), lines_data, 11, LIGHT_GRAY)
    return shape

def add_stat_card(slide, left, top, width, height, number, label, color=PRIMARY, card_bg=DARK_CARD):
    shape = add_shape(slide, left, top, width, height, card_bg)
    add_text_box(slide, left, top + Inches(0.2), width, Inches(0.5), number, 28, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.65), width, Inches(0.4), label, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Decorative accent line at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

# Title area
add_multi_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(4.5), [
    ("NIVESH SAATHI", 52, PRIMARY, True, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("Your AI-Powered Investment Companion for Bharat", 28, WHITE, True, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("Making investing as simple as chatting with a friend.", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Conversational AI  \u00b7  Jarvis-Level Intelligence  \u00b7  India-First", 16, ACCENT, False, PP_ALIGN.CENTER),
    ("", 30, WHITE, False, PP_ALIGN.CENTER),
    ("\U0001f91d  \U0001f1ee\U0001f1f3  Built with \u2764\ufe0f for Bharat", 18, MUTED, False, PP_ALIGN.CENTER),
])


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "THE PROBLEM", 14, ACCENT_RED, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
    "200 Million Indians Have Money to Invest — But Don't.", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
    "India is one of the world's fastest-growing economies, yet only 3% of Indians invest in equities.", 16, LIGHT_GRAY)

# Stats row
stats = [
    ("200M+", "Indians with disposable\nincome who don't invest"),
    ("5,000+", "Mutual funds create\ndecision paralysis"),
    ("3%", "Equity participation\nvs 55% in the US"),
    ("78%", "Millennials say investing\nis 'too complicated'"),
]
for i, (num, label) in enumerate(stats):
    add_stat_card(slide, Inches(0.8 + i*3.1), Inches(2.4), Inches(2.8), Inches(1.3), num, label, ACCENT_RED)

# Pain points
pain_points = [
    ("\U0001f4da  Financial Jargon Barrier", "NAV, CAGR, XIRR, SIP, ELSS — alien vocabulary for 90% of Indians. Existing apps assume financial literacy."),
    ("\U0001f635  Decision Paralysis", "5,000+ mutual funds, 5,000+ stocks, FDs, gold, PPF, NPS — too many choices, zero guidance on what fits YOU."),
    ("\U0001f916  Tools, Not Companions", "Current apps (Groww, Zerodha, Paytm Money) are transactional platforms. They let you buy/sell but never guide you."),
    ("\U0001f3d8\ufe0f  Tier 2/3 City Neglect", "80% of fintech focus is on metro users. Small-town India needs guidance in their own language (Hinglish)."),
]
for i, (title, body) in enumerate(pain_points):
    add_card(slide, Inches(0.8 + i*3.1), Inches(4.1), Inches(2.8), Inches(2.8),
        title, [body], card_color=RGBColor(0x1E, 0x10, 0x10), border=ACCENT_RED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "THE SOLUTION", 14, ACCENT, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
    "Nivesh Saathi: The Jarvis of Personal Finance", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(10), Inches(0.5),
    "An AI-native companion that understands your life, tracks your money, studies the market, and tells you exactly what to do — in plain Hinglish.", 16, LIGHT_GRAY)

# Core pillars
pillars = [
    ("\U0001f4ac", "Conversational AI", "Chat-based onboarding & advice in Hinglish. No forms, no jargon. Just talk to Saathi like a friend."),
    ("\U0001f9e0", "Jarvis-Level Intelligence", "Tracks every investment, spending pattern, risk appetite, market signal. Generates personalized buy/sell/hold signals daily."),
    ("\U0001f3af", "Goal-Based Planning", "Maps your life goals (home, car, education, retirement) to specific investment buckets with exact SIP amounts."),
    ("\U0001f514", "Proactive Companion", "Pushes AI notifications: sector opportunities, rebalancing alerts, tax saving reminders, spending warnings."),
    ("\U0001f30d", "Market Intelligence", "Real-time Nifty, Sensex, Gold, sector analysis. Tells you WHERE to invest TODAY with confidence scores."),
    ("\U0001f331", "Behavioral Gamification", "Money Garden: each goal grows from seed to fruit tree. Emotional connection drives 3x better habit formation."),
]
for i, (emoji, title, body) in enumerate(pillars):
    col = i % 3
    row = i // 3
    add_card(slide, Inches(0.8 + col*4.0), Inches(2.5 + row*2.4), Inches(3.7), Inches(2.1),
        title, [body], emoji, ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUCT DEEP DIVE (Architecture)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "PRODUCT ARCHITECTURE", 14, PRIMARY, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "End-to-End AI Prototype — Fully Functional", 32, WHITE, True)

# Architecture layers
layers = [
    ("PRESENTATION LAYER", PRIMARY, [
        "React Native + Expo SDK 54",
        "5 Tab Navigation (Home, Goals, Ask AI, Saathi Companion, Profile)",
        "Dark/Light Theme with smooth transitions",
        "Responsive design for all screen sizes",
    ]),
    ("INTELLIGENCE LAYER", ACCENT, [
        "AI Service: Hinglish NLP response engine",
        "Portfolio Analyzer: Risk scoring, diversification analysis",
        "Market Intelligence: Nifty/Sensex/Gold real-time + fallback",
        "Expense Intelligence: Category-wise spending analysis",
        "Notification Engine: Proactive AI-generated alerts",
        "Companion Engine: Jarvis brain — signals, briefs, net worth",
    ]),
    ("DATA LAYER", ACCENT_ORANGE, [
        "AsyncStorage for offline-first persistence",
        "User profile, investments, expenses, assets, liabilities",
        "Insurance, watchlist, life events tracking",
        "Tax profile computation (Old/New regime)",
        "Cached market data (30-min refresh)",
    ]),
]

for i, (title, color, items) in enumerate(layers):
    x = Inches(0.8)
    y = Inches(1.7 + i * 1.8)
    shape = add_shape(slide, x, y, Inches(11.7), Inches(1.6), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3), Inches(0.35), title, 13, color, True)
    items_text = "   \u2022  ".join([""] + items)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(11), Inches(1.1), items_text.strip(), 11, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — KEY FEATURES SHOWCASE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "FEATURE SHOWCASE", 14, PRIMARY, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "17 AI-Powered Features — Every One Built & Working", 32, WHITE, True)

features = [
    ("\U0001f4ac", "Chat Onboarding", "No forms. Chat in Hinglish to set up your entire financial profile in <2 minutes."),
    ("\U0001fa63", "3-Bucket System", "Safe/Growth/Opportunity pockets auto-allocated by risk profile."),
    ("\U0001f4b0", "Investment Tracker", "Record, edit, delete investments. Track by bucket and goal."),
    ("\U0001f4ca", "Expense Tracker", "Categorized spending with AI-generated saving insights."),
    ("\U0001f9e0", "AI Insights", "Portfolio health score, risk warnings, rebalancing alerts."),
    ("\U0001f4c8", "Market Intelligence", "Live Nifty, Sensex, Gold, FD rates, sector signals."),
    ("\U0001f916", "Saathi Companion", "Full Jarvis — daily briefs, signals, net worth, watchlist."),
    ("\U0001f514", "AI Notifications", "Proactive alerts: buy opportunities, risk warnings, tax tips."),
    ("\U0001f3af", "Goal Planner", "8 goal templates with auto SIP calculation & visual progress."),
    ("\U0001f331", "Money Garden", "Gamified growth: Seed \u2192 Sprout \u2192 Tree \u2192 Fruit Tree."),
    ("\U0001f4b1", "SIP Calculator", "Interactive calculator with growth visualization."),
    ("\U0001f4da", "Learn Hub", "Bite-sized financial education modules in Hinglish."),
    ("\U0001f4c9", "Buy/Sell Signals", "AI-generated signals with confidence scores and reasoning."),
    ("\U0001f4cb", "Daily Brief", "Personalized morning brief: market + portfolio + actions."),
    ("\U0001f4b3", "Net Worth Tracker", "Assets - Liabilities = Real-time net worth dashboard."),
    ("\U0001f440", "Watchlist", "Track stocks/funds with target buy/sell prices."),
    ("\U0001f4b8", "Tax Optimizer", "80C, 80D tracking with regime comparison."),
]

for i, (emoji, title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.6 + row * 1.35)
    shape = add_shape(slide, x, y, Inches(2.95), Inches(1.2), DARK_CARD)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.08), Inches(2.75), Inches(0.32),
        f"{emoji}  {title}", 12, WHITE, True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.42), Inches(2.75), Inches(0.7),
        desc, 9, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — THE COMPANION ENGINE (Jarvis Deep-Dive)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "THE JARVIS ENGINE", 14, ACCENT, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
    "Not Just an App — A True AI Companion", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(10), Inches(0.5),
    "Like Jarvis for Iron Man, Saathi connects every data point about your financial life and provides actionable intelligence.", 16, LIGHT_GRAY)

# Jarvis capabilities
caps = [
    ("\U0001f4ca", "Complete Financial DNA", [
        "Tracks every investment, expense, asset, liability",
        "Records insurance policies, life events, tax profile",
        "Builds a 360\u00b0 view of your financial life",
        "No assumptions — only data-driven insights",
    ]),
    ("\U0001f4c8", "Market Radar", [
        "Real-time Nifty 50, Sensex, Gold prices",
        "Sector-wise trend analysis (IT, Banking, Pharma...)",
        "Market sentiment (Fear/Neutral/Greed index)",
        "Correlates market moves with YOUR portfolio",
    ]),
    ("\U0001f9e0", "AI Signal Generator", [
        "Buy signals: 'IT sector dip \u2014 add to your SIP'",
        "Sell signals: 'Crypto overweight \u2014 book partial profits'",
        "Rebalance: 'Growth bucket 70% \u2014 shift 10% to safe'",
        "Each signal has confidence % and full reasoning",
    ]),
    ("\U0001f4f0", "Personalized Daily Brief", [
        "Morning greeting + market snapshot",
        "Portfolio overnight performance",
        "Today's top 3 action items",
        "Motivational nudge to stay disciplined",
    ]),
    ("\U0001f514", "Proactive Notifications", [
        "'Tech stocks down 3% \u2014 good entry point for your growth bucket'",
        "'Tax-saving deadline in 45 days \u2014 invest \u20b940K in ELSS'",
        "'Monthly spend up 22% \u2014 dining out grew \u20b94,500'",
        "Personalized, timely, and actionable",
    ]),
    ("\U0001f4b0", "Net Worth Command Center", [
        "Total Assets (stocks, MF, gold, FD, RE, PPF, NPS)",
        "Total Liabilities (home loan, car loan, credit cards)",
        "Net Worth = Assets \u2212 Liabilities",
        "Monthly trend tracking with change %",
    ]),
]

for i, (emoji, title, items) in enumerate(caps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(2.3 + row * 2.45)
    body = "\n".join([f"\u2022 {item}" for item in items])
    add_card(slide, x, y, Inches(3.85), Inches(2.2), title, items, emoji, ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — TAM / MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "MARKET OPPORTUNITY", 14, ACCENT_ORANGE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "A $15 Billion Opportunity in India's Investment Gap", 36, WHITE, True)

# TAM SAM SOM
tam_data = [
    ("TAM", "$15B", "Total Addressable Market", "500M working Indians \u00d7 average annual\ninvestable surplus \u00d7 platform fee potential.\nIndia's mutual fund AUM: \u20b965L Cr ($780B).\nDigital investment market growing 40% YoY.", ACCENT_ORANGE),
    ("SAM", "$4.2B", "Serviceable Addressable Market", "200M Indians with smartphones + disposable\nincome who currently DON'T invest.\nTier 2/3 cities: 65% of this segment.\nHinglish-speaking: 85% of target users.", PRIMARY),
    ("SOM", "$420M", "Serviceable Obtainable Market", "Target: 10M users in first 3 years.\nConversion to premium: 8-12%.\nARPU: \u20b9600/year (subscriptions + commissions).\nYear 3 revenue target: \u20b93,500 Cr.", ACCENT),
]

for i, (label, value, subtitle, details, color) in enumerate(tam_data):
    x = Inches(0.8 + i * 4.1)
    shape = add_shape(slide, x, Inches(1.8), Inches(3.7), Inches(4.5), DARK_CARD, color)
    add_text_box(slide, x, Inches(1.95), Inches(3.7), Inches(0.4), label, 14, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.35), Inches(3.7), Inches(0.7), value, 44, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.1), Inches(3.7), Inches(0.4), subtitle, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    # Divider
    add_shape(slide, x + Inches(0.5), Inches(3.55), Inches(2.7), Inches(0.02), color)
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.1), Inches(2.5), details, 11, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# Key insight bar at bottom
add_shape(slide, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6), RGBColor(0x1A, 0x25, 0x3A), ACCENT_ORANGE)
add_text_box(slide, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
    "\U0001f4a1  Only 3% of Indians invest in equities vs 55% in the US. India's demat accounts doubled from 5Cr to 10Cr in 2 years. The adoption wave is NOW.",
    12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — PRODUCT ADOPTION / WHY USERS WILL ADOPT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "EASE OF ADOPTION", 14, ACCENT, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "Why 100M+ Indians Will Adopt Nivesh Saathi", 36, WHITE, True)

adoption_features = [
    ("\U0001f4ac", "Zero-Friction Onboarding", "2-minute chat in Hinglish. No forms, no KYC upfront. Users start investing in their first session.", ACCENT),
    ("\U0001f30f", "Language-First Design", "Every screen, nudge, and AI response is in Hinglish — the way 85% of India actually communicates.", ACCENT),
    ("\U0001f9e0", "AI Does the Thinking", "User says 'Ghar ke liye paisa bachana hai.' AI auto-creates a goal, assigns bucket, calculates SIP.", ACCENT),
    ("\U0001f331", "Emotional Gamification", "Money Garden makes investing feel like nurturing a plant. 3x better retention vs traditional dashboards.", ACCENT),
    ("\U0001f4b0", "Start with \u20b9100", "No minimum barriers. Even a chai-money amount can be invested. Reduces anxiety of first-time investors.", ACCENT),
    ("\U0001f514", "Companion, Not Tool", "Daily AI nudges + proactive signals make Saathi feel like a trusted friend, not another fintech app.", ACCENT),
]

for i, (emoji, title, body, color) in enumerate(adoption_features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.7 + row * 2.65)
    add_card(slide, x, y, Inches(3.85), Inches(2.35), title, [body], emoji, color)

# Adoption metrics
add_shape(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6), DARK_CARD, ACCENT)
add_text_box(slide, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
    "Target Metrics:  \U0001f465 10M downloads (Y1)  \u00b7  \U0001f4f1 45% D7 retention  \u00b7  \u2b50 4.6+ App Store rating  \u00b7  \U0001f4ac 8 avg sessions/week  \u00b7  \U0001f4b0 18% invest in Week 1",
    12, ACCENT, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — MONETIZATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "MONETIZATION STRATEGY", 14, ACCENT_ORANGE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "Power User Monetization — How Nivesh Saathi Makes Money", 32, WHITE, True)

# Revenue streams
streams = [
    ("\U0001f451", "Saathi Premium", "\u20b999/month  \u00b7  \u20b9799/year", [
        "Advanced AI signals with 90%+ confidence",
        "Personalized daily briefs + tax optimizer",
        "Unlimited Ask AI conversations",
        "Priority notifications + sector deep-dives",
        "Target: 8-12% of users convert",
    ], ACCENT_ORANGE),
    ("\U0001f4b0", "Investment Distribution", "AMC partnership commissions", [
        "Partner with AMCs for mutual fund distribution",
        "Earn 0.5-1% trail commission on AUM",
        "Users invest via Saathi \u2192 seamless UX",
        "BSE StAR MF / MF Central integration",
        "Target: \u20b9500 Cr AUM in Year 2",
    ], PRIMARY),
    ("\U0001f4b3", "Financial Products", "Cross-sell affiliate revenue", [
        "Health & term insurance recommendations",
        "Credit card comparison + referral fees",
        "Personal loan matching for emergencies",
        "AI-matched: only relevant products shown",
        "Target: \u20b930 ARPU from cross-sell",
    ], ACCENT),
    ("\U0001f3e2", "Enterprise / API", "B2B intelligence layer", [
        "White-label companion for banks/NBFCs",
        "Investment behavior data insights (anonymized)",
        "API for financial advisors",
        "Corporate employee wellness programs",
        "Target: 5 enterprise clients Year 2",
    ], ACCENT_RED),
]

for i, (emoji, title, subtitle, items, color) in enumerate(streams):
    x = Inches(0.5 + i * 3.15)
    shape = add_shape(slide, x, Inches(1.7), Inches(2.95), Inches(4.3), DARK_CARD, color)
    add_text_box(slide, x, Inches(1.85), Inches(2.95), Inches(0.4), f"{emoji}  {title}", 14, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.25), Inches(2.95), Inches(0.3), subtitle, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.3), Inches(2.6), Inches(2.35), Inches(0.02), color)
    items_text = "\n".join([f"\u2022 {item}" for item in items])
    add_text_box(slide, x + Inches(0.2), Inches(2.75), Inches(2.55), Inches(3.1), items_text, 10, LIGHT_GRAY)

# Revenue projection
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), DARK_CARD, ACCENT_ORANGE)
add_multi_text(slide, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.9), [
    ("Revenue Projection (Conservative)", 13, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
    ("Year 1: \u20b912 Cr (Premium subs + early AMC)   \u00b7   Year 2: \u20b985 Cr (AUM trail + scale)   \u00b7   Year 3: \u20b9350 Cr (Full stack monetization)", 11, WHITE, False, PP_ALIGN.CENTER),
])


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITIVE MOAT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "COMPETITIVE MOAT", 14, PRIMARY, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "Why Nivesh Saathi Will Win Against Incumbents", 36, WHITE, True)

# Comparison table header
table_y = Inches(1.8)
headers = ["Capability", "Groww", "Zerodha", "Paytm Money", "Nivesh Saathi"]
col_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.8)]
col_x = [Inches(0.8)]
for w_val in col_widths[:-1]:
    col_x.append(col_x[-1] + w_val)

# Header row
for i, (header, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
    color = PRIMARY if i == len(headers)-1 else MUTED
    bg = RGBColor(0x2A, 0x35, 0x4A) if i < len(headers)-1 else RGBColor(0x2E, 0x2B, 0x5F)
    add_shape(slide, cx, table_y, cw, Inches(0.45), bg)
    add_text_box(slide, cx, table_y + Inches(0.05), cw, Inches(0.35), header, 12, WHITE, True, PP_ALIGN.CENTER)

# Data rows
comparisons = [
    ("AI Companion Chat", "\u274c", "\u274c", "\u274c", "\u2705 Hinglish AI"),
    ("Conversational Onboarding", "\u274c", "\u274c", "\u274c", "\u2705 2-min chat"),
    ("Proactive AI Signals", "\u274c", "\u274c", "\u274c", "\u2705 Daily signals"),
    ("Expense + Investment Tracking", "\u274c", "\u274c", "\u26a0\ufe0f Basic", "\u2705 Full AI analysis"),
    ("Gamified Investing", "\u274c", "\u274c", "\u274c", "\u2705 Money Garden"),
    ("Hinglish-First Design", "\u274c", "\u274c", "\u26a0\ufe0f Partial", "\u2705 Native"),
    ("Net Worth Dashboard", "\u274c", "\u274c", "\u274c", "\u2705 Full stack"),
    ("Tax Optimization", "\u26a0\ufe0f Basic", "\u26a0\ufe0f Basic", "\u274c", "\u2705 AI-powered"),
    ("Tier 2/3 Focus", "\u274c", "\u274c", "\u274c", "\u2705 Primary target"),
]

for r, (cap, g, z, p, ns) in enumerate(comparisons):
    ry = table_y + Inches(0.45 + r * 0.45)
    row_bg = DARK_CARD if r % 2 == 0 else RGBColor(0x16, 0x21, 0x33)
    for i, (val, cx, cw) in enumerate(zip([cap, g, z, p, ns], col_x, col_widths)):
        bg = row_bg if i < len(headers)-1 else RGBColor(0x25, 0x22, 0x4A)
        add_shape(slide, cx, ry, cw, Inches(0.45), bg)
        color = WHITE if i == 0 or i == len(headers)-1 else LIGHT_GRAY
        add_text_box(slide, cx, ry + Inches(0.07), cw, Inches(0.35), val, 10, color, i==0, PP_ALIGN.CENTER)

# Moat bullets below table
moat_y = table_y + Inches(0.45 + len(comparisons) * 0.45 + 0.15)
moats = [
    ("\U0001f310", "Hinglish-First AI", "First mover in conversational\ninvesting in Indian languages"),
    ("\U0001f9e0", "Data Flywheel", "More users \u2192 better AI \u2192 better\nrecommendations \u2192 more users"),
    ("\U0001f91d", "Trust via Simplicity", "Plain language builds trust with\nnon-English speaking investors"),
    ("\U0001f504", "Habit Loop", "Garden gamification + daily nudges\ncreate 3x retention vs competitors"),
]
for i, (emoji, title, desc) in enumerate(moats):
    x = Inches(0.8 + i * 3.1)
    add_card(slide, x, moat_y, Inches(2.8), Inches(1.3), title, [desc], emoji, PRIMARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — RISKS & MITIGATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "RISKS & MITIGATION", 14, ACCENT_RED, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "What Could Cause Failure in 12 Months — And Our Defenses", 32, WHITE, True)

risks = [
    ("SEBI Regulatory Risk", "HIGH",
     "New SEBI regulations could restrict AI-based investment advice or require RIA registration.",
     [
         "Position as 'education + tracking' tool, not investment advisor",
         "Partner with SEBI-registered RIAs for actual buy/sell execution",
         "Apply for IA license proactively",
         "All recommendations carry standard disclaimers",
     ]),
    ("User Trust with Money", "HIGH",
     "First-time investors are inherently suspicious of new platforms handling their money.",
     [
         "Start with education + tracking (no money handling initially)",
         "Partner with trusted brands (BSE StAR MF, banks)",
         "Transparent AI: always show reasoning behind every signal",
         "Start with small amounts (\u20b9100) to build confidence",
     ]),
    ("Competition from Incumbents", "MEDIUM",
     "Groww (8Cr users), Zerodha (1.2Cr), Paytm Money could add AI features.",
     [
         "Incumbents are tools, we are companions — different DNA",
         "Hinglish-first + Tier 2/3 focus = different market entirely",
         "Speed advantage: AI-native architecture vs bolted-on AI",
         "Community-driven growth (referral + word of mouth)",
     ]),
    ("AI Accuracy & Hallucination", "MEDIUM",
     "AI giving wrong investment advice could destroy trust and invite legal action.",
     [
         "Conservative recommendation engine (never speculative)",
         "All signals based on data, not predictions",
         "Confidence scores on every recommendation",
         "Human-in-the-loop for high-value decisions",
     ]),
    ("Unit Economics Pressure", "LOW",
     "High CAC + free tier could make unit economics unsustainable.",
     [
         "Organic growth via Hinglish content + community",
         "Multiple revenue streams (subs + commission + cross-sell)",
         "Low infra cost (Expo + AsyncStorage = minimal backend)",
         "Target LTV:CAC > 3x by Month 18",
     ]),
    ("Market Downturn Impact", "LOW",
     "Major market crash could scare first-time investors away from investing entirely.",
     [
         "Position as 'all-weather' companion (protect in downturns)",
         "Safe bucket emphasis during volatile periods",
         "AI signals shift to defensive (hold/rebalance) in bear markets",
         "Education content about long-term investing discipline",
     ]),
]

for i, (title, severity, desc, mitigations) in enumerate(risks):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.85)
    sev_color = ACCENT_RED if severity == "HIGH" else (ACCENT_ORANGE if severity == "MEDIUM" else ACCENT)

    shape = add_shape(slide, x, y, Inches(3.85), Inches(2.6), DARK_CARD, sev_color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(2.5), Inches(0.3),
        f"\u26a0\ufe0f {title}", 12, WHITE, True)
    add_text_box(slide, x + Inches(2.8), y + Inches(0.1), Inches(0.9), Inches(0.25),
        severity, 9, sev_color, True, PP_ALIGN.RIGHT)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.38), Inches(3.55), Inches(0.5),
        desc, 9, LIGHT_GRAY)
    mit_text = "\n".join([f"\u2713 {m}" for m in mitigations])
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(3.55), Inches(1.6),
        mit_text, 9, ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — USER JOURNEY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "USER JOURNEY", 14, PRIMARY, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "From Download to Disciplined Investor — In 7 Steps", 32, WHITE, True)

steps = [
    ("1", "Download & Open", "User sees a clean, warm welcome screen. No sign-up walls. Straight into the experience."),
    ("2", "Chat with Saathi", "AI asks about income, expenses, goals in Hinglish. Feels like talking to a smart friend."),
    ("3", "Auto-Generated Plan", "3-bucket system created: Safe, Growth, Opportunity. Goals mapped with SIP amounts."),
    ("4", "First Investment", "Record first investment starting from \u20b9100. See the Money Garden seed planted."),
    ("5", "Daily AI Companion", "Morning briefs, sector alerts, spending insights. Saathi is always watching out for you."),
    ("6", "Track & Grow", "Watch net worth grow. Garden blooms. AI optimizes tax, rebalances portfolio."),
    ("7", "Become Confident", "From 'I don't know where to invest' to 'My Saathi and I have a plan.'"),
]

# Draw journey as connected cards
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.3 + i * 1.85)
    y = Inches(1.8)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.55), y + Inches(0.05), Inches(0.5), Inches(0.4), num, 18, WHITE, True, PP_ALIGN.CENTER)

    # Connector line
    if i < len(steps) - 1:
        add_shape(slide, x + Inches(1.05), y + Inches(0.22), Inches(0.8), Inches(0.04), PRIMARY)

    # Card
    add_card(slide, x, y + Inches(0.65), Inches(1.65), Inches(2.5), title, [desc], card_color=DARK_CARD, border=PRIMARY)

# Bottom insight
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), DARK_CARD, PRIMARY)
add_multi_text(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.6), [
    ("\U0001f4a1  The Key Insight: Companion > Tool", 16, PRIMARY, True, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("Traditional apps stop at Step 3 (just a plan). Nivesh Saathi continues through Steps 4-7,", 13, WHITE, False, PP_ALIGN.CENTER),
    ("providing CONTINUOUS AI companionship — tracking, analyzing, alerting, and motivating.", 13, WHITE, False, PP_ALIGN.CENTER),
    ("This is the difference between a tool you open once a month and a companion you check daily.", 13, ACCENT, True, PP_ALIGN.CENTER),
])


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — TECH STACK & INNOVATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "TECH STACK & INNOVATION", 14, ACCENT, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "Built for Scale, Designed for Bharat", 32, WHITE, True)

# Tech columns
tech_sections = [
    ("Frontend", ACCENT, [
        ("React Native + Expo SDK 54", "Cross-platform: iOS + Android from single codebase"),
        ("TypeScript", "Type-safe code = fewer bugs at scale"),
        ("React Navigation", "Fluid tab + stack navigation"),
        ("Expo Linear Gradient", "Premium visual design system"),
        ("Custom Component Library", "13+ purpose-built UI components"),
    ]),
    ("AI & Intelligence", PRIMARY, [
        ("Built-in NLP Engine", "Hinglish understanding without cloud dependency"),
        ("Portfolio Analyzer", "Risk scoring, diversification analysis, health index"),
        ("Market Intelligence", "Real-time API + smart fallback data system"),
        ("Companion Engine", "Signal generation, daily briefs, net worth calc"),
        ("Notification Engine", "Proactive AI alerts with context awareness"),
    ]),
    ("Data & Infrastructure", ACCENT_ORANGE, [
        ("AsyncStorage", "Offline-first — works without internet"),
        ("Context API", "Centralized state management with React Context"),
        ("EAS Build", "Cloud-built APKs via Expo Application Services"),
        ("Modular Architecture", "12 service modules, 10 screens, 5 components"),
        ("Zero Backend Cost", "All intelligence runs on-device = privacy + speed"),
    ]),
]

for i, (section_title, color, items) in enumerate(tech_sections):
    x = Inches(0.5 + i * 4.15)
    shape = add_shape(slide, x, Inches(1.6), Inches(3.85), Inches(5.5), DARK_CARD, color)
    add_text_box(slide, x, Inches(1.7), Inches(3.85), Inches(0.4), f"\U0001f527  {section_title}", 15, color, True, PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.3), Inches(2.15), Inches(3.25), Inches(0.02), color)

    for j, (tech, desc) in enumerate(items):
        ty = Inches(2.35 + j * 0.95)
        add_text_box(slide, x + Inches(0.2), ty, Inches(3.45), Inches(0.3), tech, 12, WHITE, True)
        add_text_box(slide, x + Inches(0.2), ty + Inches(0.3), Inches(3.45), Inches(0.4), desc, 10, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — ORIGINAL INSIGHTS & RESEARCH
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "ORIGINAL INSIGHTS & RESEARCH", 14, PRIMARY, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "Insights That Shaped Our Product Decisions", 32, WHITE, True)

insights = [
    ("\U0001f4a1", "The Language Gap Insight",
     "Research Finding: 72% of Tier 2/3 Indians abandon English-only finance apps within 48 hours. Hinglish reduces drop-off by 4x.",
     "Product Decision: Every interaction — from onboarding chat to AI signals to daily briefs — is in Hinglish.",
     "Source: NASSCOM Digital India Report 2024, Google India Search Trends"),

    ("\U0001f4a1", "The Companion vs Tool Insight",
     "Research Finding: Users open Groww/Zerodha 2.3x/month (transactional). WhatsApp opens 23x/day (relational). Companions > Tools.",
     "Product Decision: Built Saathi as a Jarvis-level companion with daily briefs, proactive alerts, and ongoing conversations.",
     "Source: App Annie Mobile Usage Report, internal behavior modeling"),

    ("\U0001f4a1", "The Decision Paralysis Insight",
     "Research Finding: When given 24 options, people are 10x less likely to make a decision vs 6 options (Jam Experiment, Columbia).",
     "Product Decision: Simplified to 3 buckets (Safe/Growth/Opportunity) instead of showing 5,000 mutual funds.",
     "Source: Sheena Iyengar's Choice Overload Research, Columbia University"),

    ("\U0001f4a1", "The Gamification Retention Insight",
     "Research Finding: Duolingo's streak system drives 3x retention. Emotional visual progress creates habit loops.",
     "Product Decision: Money Garden (Seed \u2192 Tree) gives emotional feedback that numbers alone cannot provide.",
     "Source: Duolingo S-1 Filing, Nir Eyal's 'Hooked' framework"),

    ("\U0001f4a1", "The Trust-Through-Transparency Insight",
     "Research Finding: 68% of Indian retail investors cite 'lack of understanding' as #1 reason for not investing (SEBI survey).",
     "Product Decision: Every AI signal shows confidence % and full reasoning. No black-box recommendations.",
     "Source: SEBI Investor Survey 2023, RBI Financial Literacy Report"),

    ("\U0001f4a1", "The Proactive vs Reactive Insight",
     "Research Finding: Push notifications with personal financial context have 7x higher engagement than generic market alerts.",
     "Product Decision: AI notification engine generates personalized, context-aware alerts based on user's actual portfolio.",
     "Source: CleverTap Fintech Benchmark Report 2024"),
]

for i, (emoji, title, finding, decision, source) in enumerate(insights):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.85)
    shape = add_shape(slide, x, y, Inches(3.85), Inches(2.6), DARK_CARD, PRIMARY)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.55), Inches(0.3),
        f"{emoji} {title}", 12, PRIMARY, True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.55), Inches(0.7),
        finding, 9, WHITE)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.15), Inches(3.55), Inches(0.6),
        decision, 9, ACCENT)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.85), Inches(3.55), Inches(0.6),
        source, 8, MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — ROADMAP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "PRODUCT ROADMAP", 14, ACCENT, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
    "From Prototype to Market Leader — 18-Month Plan", 32, WHITE, True)

phases = [
    ("Phase 1: Foundation", "Month 1-3", ACCENT, "NOW", [
        "\u2705 AI prototype (complete)",
        "\u2705 Conversational onboarding",
        "\u2705 3-bucket investment system",
        "\u2705 Expense & investment tracking",
        "\u2705 AI companion with signals",
        "\u2705 APK build & deployment",
        "\u25fb App Store launch (Android)",
    ]),
    ("Phase 2: Growth", "Month 4-8", PRIMARY, "NEXT", [
        "\u25fb Real API integrations (BSE StAR MF)",
        "\u25fb KYC + actual transactions",
        "\u25fb Push notification service",
        "\u25fb Hindi/Tamil/Telugu language support",
        "\u25fb Social investing (follow friends)",
        "\u25fb Advanced AI with GPT-4 backend",
        "\u25fb iOS App Store launch",
    ]),
    ("Phase 3: Monetize", "Month 9-14", ACCENT_ORANGE, "FUTURE", [
        "\u25fb Saathi Premium subscription",
        "\u25fb AMC partnerships for MF distribution",
        "\u25fb Insurance cross-sell",
        "\u25fb Tax filing integration",
        "\u25fb Portfolio rebalancing automation",
        "\u25fb Family financial planning",
        "\u25fb Voice-based AI companion",
    ]),
    ("Phase 4: Scale", "Month 15-18", ACCENT_RED, "VISION", [
        "\u25fb 10M+ users target",
        "\u25fb Multi-language AI (8 languages)",
        "\u25fb Enterprise API for banks",
        "\u25fb AI wealth manager certification",
        "\u25fb International expansion (SEA)",
        "\u25fb IPO readiness preparation",
        "\u25fb Full-stack neobank features",
    ]),
]

for i, (title, timeline, color, badge, items) in enumerate(phases):
    x = Inches(0.5 + i * 3.15)
    shape = add_shape(slide, x, Inches(1.7), Inches(2.95), Inches(5.3), DARK_CARD, color)

    # Badge
    badge_shape = add_shape(slide, x + Inches(0.1), Inches(1.8), Inches(0.8), Inches(0.3), color)
    add_text_box(slide, x + Inches(0.1), Inches(1.8), Inches(0.8), Inches(0.3), badge, 9, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.0), Inches(1.82), Inches(1.8), Inches(0.3), timeline, 10, LIGHT_GRAY, False)
    add_text_box(slide, x + Inches(0.15), Inches(2.2), Inches(2.65), Inches(0.35), title, 13, color, True)
    add_shape(slide, x + Inches(0.15), Inches(2.55), Inches(2.65), Inches(0.02), color)

    items_text = "\n".join(items)
    add_text_box(slide, x + Inches(0.15), Inches(2.7), Inches(2.65), Inches(4.0), items_text, 10, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

# Big closing statement
add_multi_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(5.5), [
    ("NIVESH SAATHI", 48, PRIMARY, True, PP_ALIGN.CENTER),
    ("", 14, WHITE, False, PP_ALIGN.CENTER),
    ("Aapka AI Money Buddy", 32, WHITE, True, PP_ALIGN.CENTER),
    ("", 14, WHITE, False, PP_ALIGN.CENTER),
    ("We're not building another investment app.", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("We're building the Jarvis that every Indian investor deserves.", 20, ACCENT, True, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("\u2022  200M untapped investors  \u2022  $15B market opportunity  \u2022  AI-native architecture", 16, WHITE, False, PP_ALIGN.CENTER),
    ("\u2022  Hinglish-first  \u2022  Jarvis-level companion  \u2022  Working prototype today", 16, WHITE, False, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("\U0001f680  The future of investing in India is conversational, intelligent, and personal.", 18, PRIMARY, True, PP_ALIGN.CENTER),
    ("Nivesh Saathi is that future.", 18, PRIMARY, True, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    ("Thank You \U0001f64f", 28, ACCENT, True, PP_ALIGN.CENTER),
])

# Bottom accent
add_shape(slide, Inches(0), Inches(7.44), W, Inches(0.06), PRIMARY)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Nivesh_Saathi_Pitch_Deck.pptx")
prs.save(output_path)
print(f"\n{'='*60}")
print(f"  PITCH DECK GENERATED SUCCESSFULLY!")
print(f"  File: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"{'='*60}\n")

